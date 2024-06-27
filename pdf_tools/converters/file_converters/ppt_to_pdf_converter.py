import logging
import io
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF

def convert(ppt_path, pdf_path):
    try:
        logging.info(f"Starting PowerPoint to PDF conversion: {ppt_path} to {pdf_path}")
        presentation = Presentation(ppt_path)
        pdf_document = fitz.open()

        for slide_index, slide in enumerate(presentation.slides):
            # Save slide as an image
            image_stream = io.BytesIO()
            slide.shapes._spTree.blip_fill._blip.stream = image_stream
            image_stream.seek(0)

            # Convert the image stream to an image
            image = Image.open(image_stream)
            image_byte_array = io.BytesIO()
            image.save(image_byte_array, format='PNG')
            image_byte_array.seek(0)

            # Create a PDF page and insert the image
            page = pdf_document.new_page(width=image.width, height=image.height) # type: ignore
            pix = fitz.Pixmap(fitz.csRGB, image_byte_array.read(), 0)
            page.insert_image(page.rect, pixmap=pix)

            logging.info(f"Added slide {slide_index + 1} to PDF")

        pdf_document.save(pdf_path)
        logging.info(f"PowerPoint to PDF conversion completed: {pdf_path}")
    except Exception as e:
        logging.error(f"Error converting PowerPoint to PDF: {e}")
    finally:
        if 'pdf_document' in locals():
            pdf_document.close()

# Example usage
if __name__ == '__main__':
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
    ppt_path = 'example.pptx'  # Replace with your PowerPoint file path
    pdf_path = 'output.pdf'    # Replace with your desired PDF output file path
    
    convert_ppt_to_pdf(ppt_path, pdf_path) # type: ignore
