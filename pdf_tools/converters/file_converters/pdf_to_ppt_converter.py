import logging
import fitz
from pptx import Presentation
from pptx.util import Inches

def convert(pdf_path, ppt_path):
    try:
        logging.info(f"Starting PDF to PowerPoint conversion: {pdf_path} to {ppt_path}")
        pdf_document = fitz.open(pdf_path)
        presentation = Presentation()
        
        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)
            pix = page.get_pixmap() # type: ignore
            image_filename = f"temp_image_{page_number+1}.png"
            pix.save(image_filename)

            slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.add_picture(image_filename, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

            logging.info(f"Added page {page_number + 1} to the PowerPoint")

        presentation.save(ppt_path)
        logging.info(f"PDF to PowerPoint conversion completed: {ppt_path}")
    except Exception as e:
        logging.error(f"Error converting PDF to PowerPoint: {e}")
    finally:
        if 'pdf_document' in locals():
            pdf_document.close()
