import fitz  # PyMuPDF
from pptx import Presentation

def convert(pdf_path, ppt_path):
    pdf_document = fitz.open(pdf_path)
    presentation = Presentation()
    
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text = page.get_text("text")
        
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"Page {page_num + 1}"
        content.text = text
    
    presentation.save(ppt_path)
    print(f'PDF converted to PowerPoint and saved as {ppt_path}')

