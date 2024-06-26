from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def convert(ppt_path, pdf_path):
    presentation = Presentation(ppt_path)
    pdf = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter
    
    for slide_num, slide in enumerate(presentation.slides):
        if slide_num > 0:
            pdf.showPage()
        
        slide_text = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + '\n'
        
        text_object = pdf.beginText()
        text_object.setTextOrigin(40, height - 40)
        text_object.setFont("Helvetica", 12)
        
        for line in slide_text.split('\n'):
            text_object.textLine(line)
        
        pdf.drawText(text_object)
    
    pdf.save()
    print(f'PowerPoint converted to PDF and saved as {pdf_path}')
